from pptx import Presentation
import os
import pandas as pd
from pptx.util import Pt
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Cm
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def create_line_slide(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.width = Pt(1)  # Set the width to 1 point
    connector.line.color.rgb = RGBColor(0, 0, 255)  # Blue color  # Change 'red' to the desired RGB color

    # connector.top = y2
    return connector

def SubElement(parent, tagname, **kwargs):
                element = OxmlElement(tagname)
                element.attrib.update(kwargs)
                parent.append(element)
                return element

def _set_cell_border(cell, border_color="0000FF", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')


#open presentation
prs = Presentation()
# Set slide width and height for 16:9 aspect ratio
prs.slide_width = Cm(33.867)  # width for 16:9
prs.slide_height = Cm(19.05)  # height for 16:9

# Define slide layout
blank_slide_layout = prs.slide_layouts[6]  # Assuming layout 6 corresponds to a blank layout
slide = prs.slides.add_slide(blank_slide_layout)
txBox = slide.shapes.add_textbox(Cm(1.91), Cm(5.92), Cm(30.72), Cm(4.08))
tf = txBox.text_frame
tf.text = 'CIS Assessment Report'  # Change to your desired text
tf.paragraphs[0].font.name = "Arial Nova"
tf.paragraphs[0].font.size = Pt(40)  # Set font size to 14 points
tf.text_anchor = PP_ALIGN.LEFT
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.paragraphs[0].font.color.rgb = RGBColor(21, 96, 130)  # Setting the font color 

# Read data from reporte.xlsx into pandas DataFrame
df = pd.read_excel("20240418_CSI_Assessments_findings_impacts_recommendation.xlsx", sheet_name="Controls V8")

# Add slides and lines based on DataFrame rows
for index, row in df.iterrows():    
    if pd.isnull(row.iloc[1]):
        
        #agrego una slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # creo la linea de arriba de todo.
        create_line_slide(slide, Cm(0.82), Cm(0.92), Cm(32.76), Cm(0.92))

        #creo el textbox para el numero
        txBox = slide.shapes.add_textbox(Cm(0.82), Cm(1.27), Cm(1.74), Cm(1.28))
        tf = txBox.text_frame
        tf.text = str(row.iloc[0])  # Change to your desired text
        tf.paragraphs[0].font.name = "Aptos Narrow"
        tf.paragraphs[0].font.size = Pt(24)  # Set font size to 14 points
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(21, 96, 130)  # Setting the font color 
        # Set text alignment
        tf.text_anchor = PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        #creo el textbox para el titulo (ej. Data Protection)
        left = Cm(2.00)
        top = Cm(1.05)
        width = Cm(32.76)
        height = Cm(0.86)
        # Add a textbox para el row[5]
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        # update_text_of_textbox(prs, 1, tf, texto)
        tf.text = row.iloc[4]  # Change to your desired text
        tf.paragraphs[0].font.name = "Aptos Narrow"
        tf.paragraphs[0].font.size = Pt(14)  # Set font size to 14 points
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Setting the font color 

        # Agrego el texto debajo de Data Protection
        width = Cm(32.76) - left
        top = top + height
        num_lines = len(row.iloc[5]) // 100 + (len(row.iloc[5]) % 20 > 0)  # This calculates the number of lines needed
        height = Pt(num_lines * 12)  # Assuming 12 points per line, adjust accordingly
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        # update_text_of_textbox(prs, 1, tf, texto)
        tf.text = row.iloc[5]  # Change to your desired text
        tf.paragraphs[0].font.name = "Aptos Narrow"
        tf.paragraphs[0].font.size = Pt(12)  # Set font size to 14 points
        tf.paragraphs[0].font.bold = False
        # Set text alignment
        tf.text_anchor = PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.TOP
        # Set text wrapping
        tf.word_wrap = True
        tf.auto_size = True

    else:
        # si el flag es "Falso" y "no tiene criticidad"        
        if (row.iloc[11] == True) and (not pd.isnull(row.iloc[10])):        
            #si se termina la pagina cambio el top para que empiece de nuevo y creo una slide nueva.
            if top > Cm(15):
                top = Cm(0.2)
                #agrego una slide
                slide = prs.slides.add_slide(blank_slide_layout)
            # creo la linea
            top = height = top + height + Cm(0.5)
            create_line_slide(slide, Cm(2.00), top, Cm(32.76), top)
            shapes = slide.shapes
            table = shapes.add_table(1, 4, Cm(24.7), top, Cm(5), Cm(0.5)).table
            for cell in table.iter_cells():
                _set_cell_border(cell)
            #arma la tabla para poner los datos a cada item.
            table.columns[0].width = Cm(2.5)
            table.columns[1].width = Cm(2.5)
            table.columns[2].width = Cm(0.8)
            table.columns[3].width = Cm(2.3)
            # celda de dispositivos/datos/Aplicaciones -> es la primera de las 4
            cell = table.cell(0, 0)        
            cell.text = str(row.iloc[2])
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(234, 234, 234)  

            # celda de Identity/Protect/Detectar
            cell = table.cell(0, 1)
            cell.text = str(row.iloc[3])
            #add color to the cell.
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # celda de Prioridad /numero.
            cell = table.cell(0, 2)
            cell.text = str(row.iloc[9])[0]
            #add color to the cell.
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # celda de Riesgos
            cell = table.cell(0, 3)
            cell.text = str(row.iloc[10])
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            if str(row.iloc[10]) == "Critical":
                cell.fill.fore_color.rgb = RGBColor(255, 51, 0)
            elif str(row.iloc[10]) == "High":
                cell.fill.fore_color.rgb = RGBColor(228, 108,10)
            elif str(row.iloc[10]) == "Medium":
                cell.fill.fore_color.rgb = RGBColor(227, 227, 11)
            else:
                cell.fill.fore_color.rgb = RGBColor(51, 153, 51)
            
            #2da tabla debajo de la tabla.
            #creo mi_top para poder poner el top de la tabla de Unidades de Negocio
            table = shapes.add_table(1, 4, Cm(16.7), top, Cm(5), Cm(0.5)).table
            for cell in table.iter_cells():
                _set_cell_border(cell)
            #arma la tabla para poner los datos a cada item.
            table.columns[0].width = Cm(2.8)
            table.columns[1].width = Cm(2)
            table.columns[2].width = Cm(1.8)
            table.columns[3].width = Cm(1.5)
            #primera celda de la 2da tabla
            cell = table.cell(0, 0)
            cell.text = 'Medios de Pagos'
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            #segunda celda de la 2da tabla
            cell = table.cell(0, 1)
            cell.text = 'Transporte'
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            #tercera celda de la 2da tabla
            cell = table.cell(0, 2) 
            cell.text = 'Salud'
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            #cuarta celda de la 2da tabla
            cell = table.cell(0, 3)
            cell.text = 'BPS'
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.name = "Aptos Narrow"
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            left = Cm(1.80)
            #creo el textbox para el numero
            txBox = slide.shapes.add_textbox(left, top + 0.1, Cm(1.15), Cm(0.77))
            tf = txBox.text_frame
            tf.text = str(row.iloc[1])  # 2do nivel de numero
            tf.paragraphs[0].font.name = "Aptos Narrow"
            tf.paragraphs[0].font.size = Pt(11)  # Set font size to 11 points
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(21, 96, 130)  # Setting the font color 

            #creo el textbox para el titulo del 2do Nivel
            left = Cm(3)
            width = Cm(16.93)
            height = Cm(0.86)
            # Add a textbox para el row[4]
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            # update_text_of_textbox(prs, 1, tf, texto)
            tf.text = row.iloc[4]  # Change to your desired text
            tf.paragraphs[0].font.name = "Aptos Narrow"
            tf.paragraphs[0].font.size = Pt(11)  # Set font size to 12 points
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Setting the font color 

            #creo el textbox del FINDING *******
            width = Cm(29.76)
            top = top + height + Cm(0.1)
            num_lines = len(row.iloc[5]) // 100 + (len(row.iloc[5]) % 30 > 0)  # This calculates the number of lines needed
            height = Pt(num_lines * 12)  # Assuming 12 points per line, adjust accordingly
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = 'Observación: '+str(row.iloc[5])  # Change to your desired text
            # Set the height of the text box
            txBox.height = height
            tf.paragraphs[0].font.name = "Aptos Narrow"
            tf.paragraphs[0].font.size = Pt(11)  # Set font size to 100 points
            # Set text alignment
            tf.text_anchor = PP_ALIGN.LEFT
            tf.vertical_anchor = MSO_ANCHOR.TOP
            # Set text wrapping
            tf.word_wrap = True
            tf.auto_size = True

            #creo el textbox del IMPACT *******
            width = Cm(29.76)
            top = top + height + Cm(0.1)
            num_lines = len(row.iloc[6]) // 100 + (len(row.iloc[6]) % 20 > 0)  # This calculates the number of lines needed
            height = Pt(num_lines * 12)  # Assuming 12 points per line, adjust accordingly
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = 'Impacto: '+str(row.iloc[6])  # Change to your desired text
            # Set the height of the text box
            txBox.height = height
            tf.paragraphs[0].font.name = "Aptos Narrow"
            tf.paragraphs[0].font.size = Pt(12)  # Set font size to 100 points
            # Set text alignment
            tf.text_anchor = PP_ALIGN.LEFT
            tf.vertical_anchor = MSO_ANCHOR.TOP
            # Set text wrapping
            tf.word_wrap = True
            tf.auto_size = True

            #creo el textbox de la RECOMENDACION *******
            width = Cm(29.76)
            top = top + height + Cm(0.5)
            num_lines = len(row.iloc[8]) // 100 + (len(row.iloc[8]) % 20 > 0)  # This calculates the number of lines needed
            height = Pt(num_lines * 12)  # Assuming 12 points per line, adjust accordingly
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = 'Recommendation: '+str(row.iloc[8])  # Change to your desired text
            # Set the height of the text box
            txBox.height = height
            tf.paragraphs[0].font.name = "Aptos Narrow "
            tf.paragraphs[0].font.size = Pt(10)  # Set font size to 100 points
            # Set text alignment
            tf.text_anchor = PP_ALIGN.LEFT
            tf.vertical_anchor = MSO_ANCHOR.TOP
            # Set text wrapping
            tf.word_wrap = True
            tf.auto_size = True

            top = top + Cm(0.1)   # Increase top position for next iteration
        else:
            continue


#end 
blank_slide_layout = prs.slide_layouts[6]  # Assuming layout 6 corresponds to a blank layout
slide = prs.slides.add_slide(blank_slide_layout)
txBox = slide.shapes.add_textbox(Cm(1.91), Cm(5.92), Cm(30.72), Cm(4.08))
tf = txBox.text_frame
tf.text = 'The End'  # Change to your desired text
tf.paragraphs[0].font.name = "Arial Nova"
tf.paragraphs[0].font.size = Pt(40)  # Set font size to 14 points
tf.text_anchor = PP_ALIGN.LEFT
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.paragraphs[0].font.color.rgb = RGBColor(21, 96, 130)  # Setting the font color 

# Save the presentation
output_path = os.path.join("ppts", "0001_SUMMARY.pptx")
prs.save(output_path)